#!/usr/bin/env python3
"""Convert the Pokemon Starters HTML presentation to PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# monday.com brand colors
BLACK = RGBColor(0x00, 0x00, 0x00)
SURFACE = RGBColor(0x23, 0x24, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SECONDARY = RGBColor(0xC3, 0xCE, 0xD8)
MUTED = RGBColor(0xA0, 0xA0, 0xA0)
PURPLE = RGBColor(0x61, 0x64, 0xFF)
RED = RGBColor(0xFF, 0x3D, 0x57)
YELLOW = RGBColor(0xFF, 0xCB, 0x00)
GREEN = RGBColor(0x00, 0xC8, 0x75)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

BLANK_LAYOUT = prs.slide_layouts[6]


def set_slide_bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=24,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_weight=400, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold or font_weight >= 600
    p.font.name = "Poppins"
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_rect(slide, left, top, width, height, color=SURFACE, radius=0.15):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = radius
    return shape


def add_multi_text(slide, left, top, width, height, lines, font_size=18,
                   color=SECONDARY, alignment=PP_ALIGN.LEFT, bullet=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, text_color, text_bold, text_size) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if bullet and not text_bold:
            p.text = f"  {text}"
        else:
            p.text = text
        p.font.size = Pt(text_size or font_size)
        p.font.color.rgb = text_color or color
        p.font.bold = text_bold
        p.font.name = "Poppins"
        p.alignment = alignment
        p.space_after = Pt(4)
    return tf


# ── Slide 1: Cover ──
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_text(slide, 2, 2.5, 12, 2, "Pokémon Yellow:", 60, WHITE, True,
         PP_ALIGN.CENTER)
add_text(slide, 2, 4.2, 12, 1.5, "Which Starter?", 60, YELLOW, True,
         PP_ALIGN.CENTER)
add_text(slide, 2, 6, 12, 1, "A guide to Bulbasaur, Charmander & Squirtle",
         28, SECONDARY, False, PP_ALIGN.CENTER)


# ── Slide 2: Meet the Starters ──
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_text(slide, 2, 0.6, 12, 1, "Meet the Starters", 44, WHITE, True,
         PP_ALIGN.CENTER)

starters = [
    ("Bulbasaur", "Grass / Poison", GREEN,
     "The balanced choice. Strong against the first two gyms (Brock & Misty) and learns status moves early.",
     "Cerulean City — from Melanie"),
    ("Charmander", "Fire", RED,
     "The hard-mode pick. Weak against the first two gyms but becomes the powerhouse Charizard at Lv. 36.",
     "Route 24 — from Damian"),
    ("Squirtle", "Water", PURPLE,
     "The fan favorite. Beats Brock easily, gets Surf (HM03), and Blastoise tanks late-game fights.",
     "Vermilion City — from Officer Jenny"),
]

card_w = 4.0
gap = 0.4
start_x = (16 - (card_w * 3 + gap * 2)) / 2

for i, (name, ptype, type_color, desc, location) in enumerate(starters):
    x = start_x + i * (card_w + gap)
    y = 2.2
    add_rect(slide, x, y, card_w, 5.8)
    add_text(slide, x, y + 0.5, card_w, 0.8, name, 28, WHITE, True,
             PP_ALIGN.CENTER)
    # Type badge
    badge_w = 2.2
    badge_x = x + (card_w - badge_w) / 2
    r, g, b = type_color[0], type_color[1], type_color[2]
    badge = add_rect(slide, badge_x, y + 1.4, badge_w, 0.5,
                     RGBColor(r // 4, g // 4, b // 4), 0.5)
    add_text(slide, badge_x, y + 1.4, badge_w, 0.5, ptype, 14, type_color,
             True, PP_ALIGN.CENTER)
    # Description
    add_text(slide, x + 0.3, y + 2.2, card_w - 0.6, 2.0, desc, 14,
             SECONDARY, False, PP_ALIGN.CENTER)
    # Location
    add_text(slide, x + 0.3, y + 4.6, card_w - 0.6, 0.6, location, 11,
             MUTED, False, PP_ALIGN.CENTER)


# ── Slide 3: Head-to-Head ──
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_text(slide, 2, 0.4, 12, 1, "Head-to-Head", 44, WHITE, True,
         PP_ALIGN.CENTER)

# Left panel - Bulbasaur
panel_w = 6.5
left_x = 1.0
right_x = 8.5
panel_y = 1.8
panel_h = 6.5

add_rect(slide, left_x, panel_y, panel_w, panel_h)
add_text(slide, left_x + 0.5, panel_y + 0.4, 5, 0.6, "Bulbasaur", 26,
         WHITE, True)

add_multi_text(slide, left_x + 0.5, panel_y + 1.3, 5.5, 3.0, [
    ("Strengths", GREEN, True, 16),
    ("Super effective vs. Gym 1 & Gym 2", SECONDARY, False, 13),
    ("Learns Leech Seed + Sleep Powder early", SECONDARY, False, 13),
    ("Resists Electric (Lt. Surge) & Grass", SECONDARY, False, 13),
    ("Available earliest — Cerulean City", SECONDARY, False, 13),
], bullet=True)

add_multi_text(slide, left_x + 0.5, panel_y + 4.0, 5.5, 2.0, [
    ("Weaknesses", RED, True, 16),
    ("Weak to Psychic (Sabrina) & Fire (Blaine)", SECONDARY, False, 13),
    ("Venusaur's movepool is limited late-game", SECONDARY, False, 13),
], bullet=True)

# Right panel - Charmander
add_rect(slide, right_x, panel_y, panel_w, panel_h)
add_text(slide, right_x + 0.5, panel_y + 0.4, 5, 0.6, "Charmander", 26,
         WHITE, True)

add_multi_text(slide, right_x + 0.5, panel_y + 1.3, 5.5, 3.0, [
    ("Strengths", GREEN, True, 16),
    ("Charizard — high Sp. Atk & Speed", SECONDARY, False, 13),
    ("Learns Slash, Flamethrower, Fire Blast", SECONDARY, False, 13),
    ("Destroys Erika (Grass) & Lorelei (Ice)", SECONDARY, False, 13),
], bullet=True)

add_multi_text(slide, right_x + 0.5, panel_y + 3.6, 5.5, 2.5, [
    ("Weaknesses", RED, True, 16),
    ("Struggles badly vs. Brock & Misty early", SECONDARY, False, 13),
    ("Available latest — Route 24", SECONDARY, False, 13),
    ("Fire types redundant with Pikachu's coverage", SECONDARY, False, 13),
], bullet=True)


# ── Slide 4: The Verdict ──
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

add_text(slide, 2, 1.5, 12, 1, "Pick", 60, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 2, 3.0, 12, 1, "Bulbasaur", 60, GREEN, True, PP_ALIGN.CENTER)
add_text(slide, 2.5, 4.5, 11, 1.5,
         "It covers the most gym weaknesses early, learns the best utility moves, "
         "and is available first. Squirtle is a close second — but Pikachu already "
         "handles Water's strengths.",
         20, SECONDARY, False, PP_ALIGN.CENTER)

# Stats row
stats = [("5/8", "Gyms with advantage"), ("Lv. 32", "Fully evolved"),
         ("#1", "Speedrun pick")]
stat_w = 3.0
stat_gap = 0.8
stats_start = (16 - (stat_w * 3 + stat_gap * 2)) / 2

for i, (num, label) in enumerate(stats):
    sx = stats_start + i * (stat_w + stat_gap)
    add_text(slide, sx, 6.2, stat_w, 0.8, num, 40, GREEN, True,
             PP_ALIGN.CENTER)
    add_text(slide, sx, 7.1, stat_w, 0.5, label, 14, MUTED, False,
             PP_ALIGN.CENTER)


# ── Slide 5: Closing ──
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_text(slide, 2, 3.0, 12, 1.5, "Now go catch 'em all.", 60, WHITE, True,
         PP_ALIGN.CENTER)
add_text(slide, 2, 5.0, 12, 1, "Gotta pick the right one first.", 24,
         SECONDARY, False, PP_ALIGN.CENTER)


# Save
output = "/Users/eliorsi/Monday.com Dropbox/Elior Siegelwachs/Initatives/Decks/monday-presentation-v2/pokemon-starters.pptx"
prs.save(output)
print(f"Saved: {output}")
